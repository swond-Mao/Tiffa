# -*- coding: utf-8 -*-
"""theme_lib.py — pptx-template-reverse 核心库
python-pptx 主题模板构建：XML 构造器 + clone_layout 配方 + 配置驱动构建。
踩坑清单见 references/clone-recipe.md，配置格式见 SKILL.md Step 3。
依赖：python-pptx>=1.0.2, lxml
"""
import copy
import json
import re

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.parts.slide import SlideLayoutPart
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI

NS = ('xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
      'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
      'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"')
EMU = 914400


def E(v):
    """inch -> EMU"""
    return int(round(v * EMU))


_id = [5000]


def nid():
    _id[0] += 1
    return _id[0]


def run_xml(text, sz, color, bold=False, font='微软雅黑'):
    esc = text.replace('&', '&amp;').replace('<', '&lt;')
    b = ' b="1"' if bold else ''
    return (f'<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="{int(sz * 100)}"{b}>'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:latin typeface="{font}"/><a:ea typeface="{font}"/></a:rPr>'
            f'<a:t>{esc}</a:t></a:r>')


def defrpr_xml(sz, color, bold=False, font='微软雅黑'):
    b = ' b="1"' if bold else ''
    return (f'<a:defRPr sz="{int(sz * 100)}"{b}>'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:latin typeface="{font}"/><a:ea typeface="{font}"/></a:defRPr>')


def ph_sp(ph_type, idx, name, x, y, w, h, text, sz, color, bold=False,
          align='l', anchor='t', font='微软雅黑'):
    """placeholder sp（带默认样式文本）。idx 编号约定见 references/ph-conventions.md"""
    rid = nid()
    algn = {'l': '', 'c': ' algn="ctr"', 'ctr': ' algn="ctr"', 'r': ' algn="r"'}[align]
    anchor = {'t': 't', 'ctr': 'ctr', 'b': 'b'}[anchor]
    body = ''
    if text:
        body = (f'<a:p><a:pPr{algn}/>{run_xml(text, sz, color, bold, font)}'
                f'<a:endParaRPr lang="zh-CN" sz="{int(sz * 100)}"/></a:p>')
    p = (f'<p:sp {NS}>'
         f'<p:nvSpPr><p:cNvPr id="{rid}" name="{name}"/>'
         f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
         f'<p:nvPr><p:ph type="{ph_type}" idx="{idx}"/></p:nvPr></p:nvSpPr>'
         f'<p:spPr><a:xfrm><a:off x="{E(x)}" y="{E(y)}"/>'
         f'<a:ext cx="{E(w)}" cy="{E(h)}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
         f'<p:txBody><a:bodyPr wrap="square" anchor="{anchor}"/>'
         f'<a:lstStyle/>{body}</p:txBody></p:sp>')
    return parse_xml(p)


def rect_sp(x, y, w, h, fill, alpha=None, line=None, line_w=0.75,
            prst='rect', adj=None, name='Deco'):
    """普通形状（rect/roundRect/ellipse/chevron…），alpha 0-100（100=不透明）"""
    rid = nid()
    geom = f'<a:prstGeom prst="{prst}"><a:avLst>{f"<a:gd name=\"adj\" fmla=\"val {adj}\"/>" if adj else ""}</a:avLst></a:prstGeom>'
    fill_xml = ''
    if fill:
        a = f'<a:alpha val="{int(alpha * 1000)}"/>' if alpha is not None and alpha < 100 else ''
        fill_xml = f'<a:solidFill><a:srgbClr val="{fill}">{a}</a:srgbClr></a:solidFill>'
    ln = '<a:ln><a:noFill/></a:ln>'
    if line:
        ln = (f'<a:ln w="{int(line_w * 12700)}"><a:solidFill>'
              f'<a:srgbClr val="{line}"/></a:solidFill></a:ln>')
    p = (f'<p:sp {NS}><p:nvSpPr><p:cNvPr id="{rid}" name="{name}"/>'
         f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
         f'<p:spPr><a:xfrm><a:off x="{E(x)}" y="{E(y)}"/>'
         f'<a:ext cx="{E(w)}" cy="{E(h)}"/></a:xfrm>{geom}{fill_xml}{ln}</p:spPr></p:sp>')
    return parse_xml(p)


def line_sp(x, y, w, fill, h=0.008, name='Line'):
    """细横线（实现=细 rect，WPS 渲染稳定）"""
    return rect_sp(x, y, w, h, fill, name=name)


def pic_sp(rId, x, y, w, h, name='Picture'):
    rid = nid()
    p = (f'<p:pic {NS}><p:nvPicPr><p:cNvPr id="{rid}" name="{name}"/>'
         f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
         f'<p:blipFill><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
         f'<p:spPr><a:xfrm><a:off x="{E(x)}" y="{E(y)}"/>'
         f'<a:ext cx="{E(w)}" cy="{E(h)}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>')
    return parse_xml(p)


def clone_layout(p, master, src_part, name):
    """克隆版式 part。五条红线（详见 references/clone-recipe.md）：
    1. partname 扫描全 package iter_parts()（含新建 part），max+1
    2. image rel 用 get_or_add 复用同一 image part（不新建）
    3. 克隆 XML 的 r:embed 按 rId_map 重写
    4. master 的 sldLayoutIdLst 必须 append p:sldLayoutId
    5. sldLayoutId id 全局唯一（跨所有 master）
    """
    used = set()
    for part in p.part.package.iter_parts():
        mm = re.match(r'/ppt/slideLayouts/slideLayout(\d+)\.xml$', str(part.partname))
        if mm:
            used.add(int(mm.group(1)))
    i = max(used) + 1 if used else 1
    pn = f'/ppt/slideLayouts/slideLayout{i}.xml'
    new_el = copy.deepcopy(src_part._element)
    new_part = SlideLayoutPart(PackURI(pn), src_part.content_type, p.part.package, new_el)
    # 复制图片 rel（复用同一 image part），记录 源rId→新rId
    rId_map = {}
    for rId, rel in src_part.rels.items():
        if rel.reltype == RT.IMAGE:
            rId_map[rId] = new_part.rels.get_or_add(rel.reltype, rel.target_part)
    for el2 in new_el.iter():
        emb = el2.get(qn('r:embed'))
        if emb is not None and emb in rId_map:
            el2.set(qn('r:embed'), rId_map[emb])
    rId = master.part.relate_to(new_part, RT.SLIDE_LAYOUT)
    sldLayoutIdLst = master._element.get_or_add_sldLayoutIdLst()
    all_ids = []
    for m in p.slide_masters:
        for e in m._element.get_or_add_sldLayoutIdLst().findall(qn('p:sldLayoutId')):
            all_ids.append(int(e.get('id')))
    new_id = max(all_ids) + 1 if all_ids else 2147483649
    sldLayoutIdLst.append(parse_xml(f'<p:sldLayoutId {NS} id="{new_id}" r:id="{rId}"/>'))
    return new_part
def find_image_part(p, filename):
    """从 package 找 media 文件名对应的 image part。
    iter_parts() 只含已物化 part；slide 引用的 image 需先物化（见 build()）。
    兜底：遍历已物化 part 的 rels 物化目标再匹配。
    """
    pkg = p.part.package

    def scan(parts):
        for part in parts:
            pn = str(part.partname)
            if pn.endswith('/' + filename) and part.content_type.startswith('image/'):
                return part
        return None

    found = scan(pkg.iter_parts())
    if found is not None:
        return found
    for part in list(pkg.iter_parts()):
        for rel in part.rels.values():
            if rel.is_external:
                continue
            t = rel.target_part
            if str(t.partname).endswith('/' + filename) and t.content_type.startswith('image/'):
                return t
    raise FileNotFoundError(f'image part not found: {filename}（确认源模板 media 里有该文件）')

def _sp_tree(layout_el):
    return layout_el.find(qn('p:cSld')).find(qn('p:spTree'))


def _c_sld(layout_el):
    return layout_el.find(qn('p:cSld'))


def apply_op(p, master, op):
    """执行单个 layout_op（rename/modify/clone），返回新 part（clone 时）"""
    type_ = op['type']
    if type_ == 'rename':
        for l in master.slide_layouts:
            if l.name == op['name']:
                l.name = op['to']
                return None
        raise KeyError(f'layout not found: {op["name"]}')

    if type_ == 'modify':
        layout = next((l for l in master.slide_layouts if l.name == op['name']), None)
        if layout is None:
            raise KeyError(f'layout not found: {op["name"]}')
        return _apply_layout_changes(p, master, layout, op)

    if type_ == 'clone':
        src = next((l for l in master.slide_layouts if l.name == op['source']), None)
        if src is None:
            raise KeyError(f'source layout not found: {op["source"]}')
        new_part = clone_layout(p, master, src.part, op['name'])
        _c_sld(new_part._element).set('name', op['name'])
        # clone 后重新获取 layout 对象视图（新 part 已挂到 master）
        layout = next(l for l in master.slide_layouts if l.name == op['name'])
        return _apply_layout_changes(p, master, layout, op)

    raise ValueError(f'unknown op type: {type_}')


def _apply_layout_changes(p, master, layout, op):
    spT = _sp_tree(layout._element)
    # 移除占位符
    for ph_type in op.get('remove_ph_types', []):
        for sp in list(spT.findall(qn('p:sp'))):
            ph = sp.find('.//' + qn('p:ph'))
            if ph is not None and ph.get('type') == ph_type:
                spT.remove(sp)
    # 追加占位符
    if op.get('add_title_ph'):
        spec = op['add_title_ph']
        spT.append(_make_ph('title', spec))
    if op.get('add_body_ph'):
        spec = op['add_body_ph']
        spT.append(_make_ph('body', spec))
    for spec in op.get('add_ph', []):
        spT.append(_make_ph('body', spec))
    # 追加装饰 shape
    for sh in op.get('add_shapes', []):
        spT.append(_make_shape(p, master, layout, sh))
    # 插入背景图（index=2 = nvGrpSpPr/grpSpPr 之后，最底层）
    if op.get('insert_bg'):
        sh = op['insert_bg']
        el = _make_shape(p, master, layout, sh)
        spT.insert(2, el)
    return None


def _make_ph(ph_type, spec):
    idx, name, x, y, w, h, text, sz, color = spec[:9]
    bold = spec[9] if len(spec) > 9 and spec[9] else False
    align = spec[10] if len(spec) > 10 and spec[10] else 'l'
    anchor = spec[11] if len(spec) > 11 and spec[11] else 't'
    return ph_sp(ph_type, int(idx), name, x, y, w, h, text, sz, color,
                 bold=bold, align=align, anchor=anchor)


def _make_shape(p, master, layout, sh):
    kind = sh['kind']
    if kind == 'rect':
        return rect_sp(sh['x'], sh['y'], sh['w'], sh['h'], sh.get('fill'),
                       alpha=sh.get('alpha'), line=sh.get('line'),
                       line_w=sh.get('line_w', 0.75), prst=sh.get('prst', 'rect'),
                       adj=sh.get('adj'), name=sh.get('name', 'Deco'))
    if kind == 'line':
        return line_sp(sh['x'], sh['y'], sh['w'], sh['fill'],
                       h=sh.get('h', 0.008), name=sh.get('name', 'Line'))
    if kind == 'pic':
        img_part = find_image_part(p, sh['source'])
        rId = layout.part.relate_to(img_part, RT.IMAGE)
        return pic_sp(rId, sh['x'], sh['y'], sh['w'], sh['h'], name=sh.get('name', 'Picture'))
    raise ValueError(f'unknown shape kind: {kind}')


def build(cfg):
    """按配置构建主题模板，返回输出路径"""
    p = Presentation(cfg['source'])
    master = p.slide_masters[cfg.get('master_index', 0)]
    # 1. 先执行 layout_ops：pic 形状通过 relate_to 把 image part 重新锚定到版式。
    #    顺序关键：iter_parts() 是图遍历，若先删示例 slide，slide 独占的 image part
    #    会失去可达性、从 iter_parts 消失，find_image_part 找不到（part 对象虽在但不会被保存）。
    for op in cfg.get('layout_ops', []):
        apply_op(p, master, op)
        print('op done:', op.get('type'), op.get('name') or op.get('to'))
    # 2. 删除示例 slide（layout_ops 完成后 image 已锚定，删页安全）
    if cfg.get('delete_sample_slides', False):
        sldIdLst = p.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn('r:id'))
            p.part.drop_rel(rId)
            sldIdLst.remove(sldId)
        print('slides removed')
    # 3. 保存
    p.save(cfg['output'])
    print('saved', cfg['output'])
    # 4. 重载自检
    p2 = Presentation(cfg['output'])
    for m in p2.slide_masters:
        for l in m.slide_layouts:
            phs = sorted(ph.placeholder_format.idx for ph in l.placeholders)
            print(' ', l.name, 'phs=', phs)
    return cfg['output']


def main():
    import sys
    if len(sys.argv) < 2:
        print('usage: python build_theme.py <theme_config.json>')
        sys.exit(1)
    cfg = json.load(open(sys.argv[1], encoding='utf-8'))
    build(cfg)


if __name__ == '__main__':
    main()
