# -*- coding: utf-8 -*-
"""probe_template.py — 设计系统提取
用法: python probe_template.py <模板.pptx>
输出：画布 / 字体 / 色板 / 版式列表(名称+占位符) / 示例页 shape 明细 / media 清单
"""
import re
import sys
from collections import Counter

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def sp_fills(sp_el):
    """收集 shape 内 solidFill 颜色（srgb/scheme）+ 线条颜色"""
    out = []
    for sf in sp_el.findall('.//{%s}solidFill' % A):
        for c in sf:
            tag = etree.QName(c).localname
            if tag == 'srgbClr':
                out.append('srgb:' + c.get('val'))
            elif tag == 'schemeClr':
                out.append('scheme:' + c.get('val'))
    for ln in sp_el.findall('.//{%s}ln' % A):
        for sf in ln.findall('{%s}solidFill' % A):
            for c in sf:
                if etree.QName(c).localname == 'srgbClr':
                    out.append('ln:' + c.get('val'))
    return out


def main():
    if len(sys.argv) < 2:
        print('usage: python probe_template.py <template.pptx>')
        sys.exit(1)
    path = sys.argv[1]
    p = Presentation(path)
    print(f'== 画布: {p.slide_width / 914400:.3f} x {p.slide_height / 914400:.3f} in')
    print(f'== slides: {len(p.slides._sldIdLst)}  masters: {len(p.slide_masters)}')

    # 字体（全 package XML 扫描 typeface）
    fonts = Counter()
    for part in p.part.package.iter_parts():
        if str(part.partname).endswith('.xml'):
            try:
                blob = part.blob
            except Exception:
                continue
            for m in re.finditer(rb'typeface="([^"]+)"', blob):
                t = m.group(1).decode('utf-8', 'ignore')
                if t and not t.startswith('+') and 'typeface' not in t:
                    fonts[t] += 1
    print('== 字体:', dict(fonts.most_common(8)))

    # 色板（全 package 扫描 srgbClr 频次，去噪：只留出现≥2次的）
    colors = Counter()
    for part in p.part.package.iter_parts():
        if str(part.partname).endswith('.xml'):
            try:
                blob = part.blob
            except Exception:
                continue
            for m in re.finditer(rb'srgbClr val="([0-9A-Fa-f]{6})"', blob):
                colors[m.group(1).decode().upper()] += 1
    top = [c for c, n in colors.most_common(30) if n >= 2]
    print('== 色板(频次>=2):', ', '.join(f'{c}×{colors[c]}' for c in top))

    # 版式列表
    for mi, m in enumerate(p.slide_masters):
        print(f'== MASTER{mi}:')
        for l in m.slide_layouts:
            phs = []
            for sp in l._element.find(qn('p:cSld')).find(qn('p:spTree')).findall(qn('p:sp')):
                ph = sp.find('.//' + qn('p:ph'))
                if ph is None:
                    continue
                off = sp.find('.//' + qn('a:off'))
                ext = sp.find('.//' + qn('a:ext'))
                pos = ''
                if off is not None and ext is not None:
                    pos = f"@({round(int(off.get('x')) / 914400, 2)},{round(int(off.get('y')) / 914400, 2)}) {round(int(ext.get('cx')) / 914400, 2)}x{round(int(ext.get('cy')) / 914400, 2)}"
                txt = sp.find('.//' + qn('a:t'))
                default = (txt.text or '')[:12] if txt is not None and txt.text else ''
                sz_el = sp.find('.//' + qn('a:defRPr'))
                sz = round(int(sz_el.get('sz')) / 100, 1) if sz_el is not None and sz_el.get('sz') else '?'
                phs.append(f"  ph idx={ph.get('idx')} type={ph.get('type')} {pos} sz={sz} default='{default}'")
            print(f'  [{l.name}]')
            for ph_line in phs:
                print(ph_line)

    # 示例页 shape 明细（前 5 页）
    for si, s in enumerate(list(p.slides)[:5]):
        print(f'== SLIDE{si + 1} (layout={s.slide_layout.name}):')
        for sh in s.shapes:
            try:
                fills = sp_fills(sh._element)
            except Exception:
                fills = []
            kind = 'PIC' if sh.shape_type == 13 else str(sh.shape_type)
            txt = ''
            if sh.has_text_frame:
                txt = sh.text_frame.text.strip().replace('\n', '|')[:30]
            print(f"  {kind} {sh.name} @({round(sh.left / 914400, 2)},{round(sh.top / 914400, 2)}) "
                  f"{round(sh.width / 914400, 2)}x{round(sh.height / 914400, 2)} "
                  f"fills={fills[:4]} text='{txt}'")

    # media 清单
    print('== MEDIA:')
    for part in p.part.package.iter_parts():
        if '/media/' in str(part.partname):
            print(f'  {part.partname} {len(part.blob)}B')


if __name__ == '__main__':
    main()
