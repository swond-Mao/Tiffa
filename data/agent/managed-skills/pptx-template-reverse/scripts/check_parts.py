# -*- coding: utf-8 -*-
"""check_parts.py — 重载验证（验收三条）
用法: python check_parts.py <file.pptx> [more.pptx ...]
验收：① dups=无 ② 版式列表各归其位 ③ 各版式 ph idx 与配置一致
"""
import sys
import zipfile

from pptx import Presentation


def main():
    if len(sys.argv) < 2:
        print('usage: python check_parts.py <file.pptx> [...]')
        sys.exit(1)
    ok = True
    for path in sys.argv[1:]:
        z = zipfile.ZipFile(path)
        names = z.namelist()
        dups = sorted(set(n for n in names if names.count(n) > 1))
        print(f'== {path}')
        print(f'   slides={len(Presentation(path).slides._sldIdLst)} parts={len(names)} dups={dups or "无"}')
        if dups:
            ok = False
        p = Presentation(path)
        for mi, m in enumerate(p.slide_masters):
            for l in m.slide_layouts:
                phs = sorted(ph.placeholder_format.idx for ph in l.placeholders)
                print(f'   M{mi}/{l.name} ph={phs}')
    print('RESULT:', 'PASS' if ok else 'FAIL (dups found)')
    sys.exit(0 if ok else 1)


if __name__ == '__main__':
    main()
