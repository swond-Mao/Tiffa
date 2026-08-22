# -*- coding: utf-8 -*-
"""qa_render.py — PIL 自渲染 deck 联系图（近似渲染，查溢出/碰撞）
用法: python qa_render.py <deck.pptx> -o sheet.png [--cols 4] [--font C:/Windows/Fonts/msyh.ttc]
版式装饰按版式自身 shape 泛化渲染（fill/line/文字/图片），不依赖具体模板。
"""
import argparse
import os
from io import BytesIO

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

S = 128.0  # px per inch


def emu2px(v):
    return v / 914400.0 * S


class Ctx:
    def __init__(self, font_path, bold_path):
        self.font_path = font_path
        self.bold_path = bold_path or font_path
        self.cache = {}

    def font(self, pt, bold=False):
        key = (round(pt), bold)
        if key not in self.cache:
            path = self.bold_path if bold else self.font_path
            self.cache[key] = ImageFont.truetype(path, max(6, int(pt * S / 72)))
        return self.cache[key]


def wrap_lines(runs, box_w, ctx):
    """runs: [(text, pt, bold, color)] -> 行列表"""
    if not runs:
        return []
    lines, cur = [], []
    cur_w = 0.0
    for run in runs:
        t, pt, bold_, color = run[0], run[1], run[2], run[3]
        f = ctx.font(pt, bold_)
        for ch in t:
            cw = f.getlength(ch)
            if cur_w + cw > box_w and cur:
                lines.append((''.join(cur), pt, bold_, color))
                cur, cur_w = [], 0.0
            cur.append(ch)
            cur_w += cw
    if cur:
        lines.append((''.join(cur), pt, bold_, color))
    return lines


def get_fill_alpha(sp_el):
    """返回 (rgb_hex, alpha 0-100) 或 (None, 100)"""
    sf = sp_el.find('.//' + qn('a:solidFill'))
    if sf is None:
        return None, 100
    srgb = sf.find(qn('a:srgbClr'))
    if srgb is None:
        return None, 100
    col = srgb.get('val')
    a_el = srgb.find(qn('a:alpha'))
    alpha = int(a_el.get('val')) / 1000 if a_el is not None else 100
    return col, alpha


def get_prst(sp_el):
    pg = sp_el.find('.//' + qn('a:prstGeom'))
    return pg.get('prst') if pg is not None else 'rect'


def shape_geom(sp_el):
    xfrm = sp_el.find('.//' + qn('a:xfrm'))
    if xfrm is None:
        return None
    off, ext = xfrm.find(qn('a:off')), xfrm.find(qn('a:ext'))
    if off is None or ext is None:
        return None
    return (emu2px(int(off.get('x'))), emu2px(int(off.get('y'))),
            emu2px(int(ext.get('cx'))), emu2px(int(ext.get('cy'))))


def shape_runs(sp_el):
    """sp 内所有段落 -> [(runs, align, sb, lsp)]"""
    out = []
    for para in sp_el.findall('.//' + qn('a:p')):
        runs = []
        for r in para.findall(qn('a:r')):
            rPr = r.find(qn('a:rPr'))
            t_el = r.find(qn('a:t'))
            if t_el is None or not (t_el.text or '').strip():
                continue
            sz = int(rPr.get('sz')) / 100 if rPr is not None and rPr.get('sz') else 10.5
            bold = rPr is not None and rPr.get('b') == '1'
            col = '4A5568'
            if rPr is not None:
                srgb = rPr.find('.//' + qn('a:srgbClr'))
                if srgb is not None:
                    col = srgb.get('val')
            runs.append((t_el.text, sz, bold, col))
        if not runs:
            continue
        al = para.find(qn('a:pPr'))
        align = 'l'
        if al is not None and al.get('algn'):
            align = {'l': 'l', 'r': 'r', 'ctr': 'ctr'}.get(al.get('algn'), 'l')
        sb = 0
        if al is not None:
            sb_el = al.find(qn('a:spcBef'))
            if sb_el is not None:
                pts = sb_el.find(qn('a:spcPts'))
                if pts is not None:
                    sb = int(pts.get('val')) / 100
        lsp = 100
        if al is not None:
            ls_el = al.find(qn('a:lnSpc'))
            if ls_el is not None:
                pct = ls_el.find(qn('a:spcPct'))
                if pct is not None:
                    lsp = int(pct.get('val')) / 1000
        out.append((runs, align, sb, lsp))
    return out


def anchor_of(sp_el):
    bodyPr = sp_el.find('.//' + qn('a:bodyPr'))
    if bodyPr is not None:
        return bodyPr.get('anchor', 't')
    return 't'


def draw_paras(img, d, x, y, w, h, paras, anchor, ctx):
    if not paras:
        return
    tot = 0.0
    for (runs, align, sb, lsp) in paras:
        for (line, s2, b2, col) in wrap_lines(runs, w, ctx):
            tot += s2 * S / 72 * 1.25 * (lsp / 100)
        tot += sb * S / 72
    yy = y + max(0, (h - tot) / 2) if anchor == 'ctr' else y
    for (runs, align, sb, lsp) in paras:
        yy += sb * S / 72
        for (line, s2, b2, col) in wrap_lines(runs, w, ctx):
            f = ctx.font(s2, b2)
            tw = f.getlength(line)
            lx = x + max(0, (w - tw) / 2) if align == 'ctr' else x + max(0, w - tw) if align == 'r' else x
            d.text((lx, yy), line, font=f, fill='#' + col)
            yy += s2 * S / 72 * 1.25 * (lsp / 100)


def draw_sp(img, d, sp_el, get_pic, ctx):
    """画单个 sp/pic（fill/line/文本/图片）"""
    geom = shape_geom(sp_el)
    if geom is None:
        return
    x, y, w, h = geom
    if w < 1 or h < 1:
        return
    # 图片（p:pic 或 sp 内嵌 blip）
    blip = sp_el.find('.//' + qn('a:blip'))
    if blip is not None:
        r_embed = blip.get(qn('r:embed'))
        pic = get_pic(r_embed) if r_embed else None
        if pic is not None:
            try:
                pim = pic.convert('RGBA').resize((max(1, int(w)), max(1, int(h))))
                img.alpha_composite(pim, (int(x), int(y)))
                d = ImageDraw.Draw(img)
            except Exception:
                d.rectangle([x, y, x + w, y + h], outline='#FF00FF')
        else:
            d.rectangle([x, y, x + w, y + h], outline='#FF00FF')
        return
    fill, alpha = get_fill_alpha(sp_el)
    prst = get_prst(sp_el)
    radius = min(w, h) * 0.12 if prst in ('roundRect', 'roundSameRectCorner') else 0
    if prst in ('ellipse', 'ellipseRibbon'):
        if fill:
            a = int(alpha * 2.55)
            if a >= 250:
                d.ellipse([x, y, x + w, y + h], fill='#' + fill)
            else:
                ov = Image.new('RGBA', img.size, (0, 0, 0, 0))
                ImageDraw.Draw(ov).ellipse([x, y, x + w, y + h],
                                           fill=(int(fill[0:2], 16), int(fill[2:4], 16), int(fill[4:6], 16), a))
                img.alpha_composite(ov)
                d = ImageDraw.Draw(img)
        return
    if fill:
        a = int(alpha * 2.55)
        if a >= 250:
            d.rounded_rectangle([x, y, x + w, y + h], radius=radius, fill='#' + fill)
        else:
            ov = Image.new('RGBA', img.size, (0, 0, 0, 0))
            od = ImageDraw.Draw(ov)
            od.rounded_rectangle([x, y, x + w, y + h], radius=radius,
                                 fill=(int(fill[0:2], 16), int(fill[2:4], 16), int(fill[4:6], 16), a))
            img.alpha_composite(ov)
            d = ImageDraw.Draw(img)
    ln = sp_el.find('.//' + qn('a:ln'))
    if ln is not None:
        srgb = ln.find('.//' + qn('a:srgbClr'))
        if srgb is not None:
            d.rounded_rectangle([x, y, x + w, y + h], radius=radius, outline='#' + srgb.get('val'))
    paras = shape_runs(sp_el)
    if paras:
        draw_paras(img, d, x, y, w, h, paras, anchor_of(sp_el), ctx)


def _pic_loader(rels):
    def get_pic(r_embed):
        if r_embed in rels:
            try:
                return Image.open(BytesIO(rels[r_embed].target_part.blob))
            except Exception:
                return None
        return None
    return get_pic


def render_layout_deco(img, d, layout, ctx):
    """按版式自身 shape 渲染装饰（含 p:pic 背景图；占位符只画默认文本提示）"""
    spTree = layout._element.find(qn('p:cSld')).find(qn('p:spTree'))
    get_pic = _pic_loader(layout.part.rels)
    for el in spTree:
        tag = el.tag.split('}')[-1]
        if tag == 'pic':
            draw_sp(img, d, el, get_pic, ctx)
            d = ImageDraw.Draw(img)
        elif tag == 'sp':
            ph = el.find('.//' + qn('p:ph'))
            if ph is not None:
                paras = shape_runs(el)
                if paras:
                    geom = shape_geom(el)
                    if geom:
                        draw_paras(img, d, geom[0], geom[1], geom[2], geom[3],
                                   paras, anchor_of(el), ctx)
                continue
            draw_sp(img, d, el, get_pic, ctx)
            d = ImageDraw.Draw(img)


def render_chart(img, d, shape, ctx):
    try:
        ch = shape.chart
        plot = ch.plots[0]
        x, y = emu2px(shape.left), emu2px(shape.top)
        w, h = emu2px(shape.width), emu2px(shape.height)
        cx, cy = x + w / 2, y + (h - 0.35 * S) / 2
        R = min(w, h - 0.35 * S) / 2 - 8
        r2 = R * 0.55
        vals = list(plot.series[0].values)
        cols = ['#0E3F8C', '#3D7BD9', '#F5A623', '#8B97A8', '#1B7A9E']
        ang = -90
        for i, v in enumerate(vals):
            sweep = v / sum(vals) * 360
            d.pieslice([cx - R, cy - R, cx + R, cy + R], ang, ang + sweep, fill=cols[i % 5])
            ang += sweep
        d.ellipse([cx - r2, cy - r2, cx + r2, cy + r2], fill='#FFFFFF')
    except Exception as e:
        d.text((x, y), f'chart? {e}', fill='red')


def render_slide(img, d, slide, ctx):
    layout = slide.slide_layout
    render_layout_deco(img, d, layout, ctx)
    d = ImageDraw.Draw(img)
    get_pic = _pic_loader(slide.part.rels)
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            render_chart(img, d, shape, ctx)
            continue
        draw_sp(img, d, shape._element, get_pic, ctx)
        d = ImageDraw.Draw(img)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('deck')
    ap.add_argument('-o', default='qa_sheet.png')
    ap.add_argument('--cols', type=int, default=4)
    ap.add_argument('--font', default=None)
    ap.add_argument('--font-bold', default=None)
    args = ap.parse_args()

    reg = args.font or 'C:/Windows/Fonts/msyh.ttc'
    bold = args.font_bold or 'C:/Windows/Fonts/msyhbd.ttc'
    if not os.path.exists(reg):
        raise SystemExit(f'font not found: {reg}（用 --font 指定）')
    if not os.path.exists(bold) and os.path.exists(reg):
        bold = reg
    ctx = Ctx(reg, bold)

    p = Presentation(args.deck)
    W = int(p.slide_width / 914400 * S)
    H = int(p.slide_height / 914400 * S)
    slides = list(p.slides)
    rows = (len(slides) + args.cols - 1) // args.cols
    tw, th = W // 2, H // 2
    sheet = Image.new('RGB', (args.cols * tw, rows * th), '#DDDDDD')
    for i, slide in enumerate(slides):
        img = Image.new('RGBA', (W, H), (255, 255, 255, 255))
        d = ImageDraw.Draw(img)
        render_slide(img, d, slide, ctx)
        im = img.convert('RGB').resize((tw, th))
        r, c = divmod(i, args.cols)
        sheet.paste(im, (c * tw, r * th))
        dd = ImageDraw.Draw(sheet)
        dd.text((c * tw + 6, r * th + 6), f'P{i + 1}',
                font=ctx.font(12, True), fill='#FF3333')
    sheet.save(args.o, quality=90)
    print('saved', args.o, sheet.size)


if __name__ == '__main__':
    main()
